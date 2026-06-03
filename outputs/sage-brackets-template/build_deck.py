"""
Sage Brackets — minimalist presentation template
Inspired by hairline technical-drawing aesthetic: corner brackets, dotted lines,
single hero color (sage), pure ink strokes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ─────────────────────────────────────────────────────────────────────────────
# DESIGN TOKENS
# ─────────────────────────────────────────────────────────────────────────────
SAGE       = RGBColor(0xD5, 0xE8, 0xB8)   # background
SAGE_DARK  = RGBColor(0xB8, 0xD0, 0x98)   # subtle fills
INK        = RGBColor(0x0F, 0x0F, 0x0F)   # strokes & primary text
INK_SOFT   = RGBColor(0x3A, 0x3A, 0x3A)   # secondary text
CLAY       = RGBColor(0xC9, 0x7B, 0x5F)   # single reserved accent
INK_MUTE   = RGBColor(0x6B, 0x6B, 0x6B)   # tertiary text / captions

FONT_DISPLAY = "Inter"          # falls back to system sans on machines without Inter
FONT_MONO    = "Consolas"       # technical-drawing label feel
FONT_BODY    = "Inter"

# ─────────────────────────────────────────────────────────────────────────────
# PRESENTATION SETUP — wide 16:9 (13.333" × 7.5")
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

# ─────────────────────────────────────────────────────────────────────────────
# LOW-LEVEL HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = SAGE
    return s


def add_line(s, x1, y1, x2, y2, color=INK, width_pt=1.0, dash=None):
    """dash: None | 'sysDot' | 'dash' | 'sysDash' etc."""
    line = s.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    if dash:
        ln = line.line._get_or_add_ln()
        # remove existing prstDash if any
        for el in ln.findall(qn('a:prstDash')):
            ln.remove(el)
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    return line


def add_corner_brackets(s, x, y, w, h, size, color=INK, width_pt=1.0):
    """Draw L-brackets at all 4 corners of an (x, y, w, h) box."""
    # Top-left
    add_line(s, x, y, x + size, y, color, width_pt)
    add_line(s, x, y, x, y + size, color, width_pt)
    # Top-right
    add_line(s, x + w - size, y, x + w, y, color, width_pt)
    add_line(s, x + w, y, x + w, y + size, color, width_pt)
    # Bottom-left
    add_line(s, x, y + h - size, x, y + h, color, width_pt)
    add_line(s, x, y + h, x + size, y + h, color, width_pt)
    # Bottom-right
    add_line(s, x + w - size, y + h, x + w, y + h, color, width_pt)
    add_line(s, x + w, y + h - size, x + w, y + h, color, width_pt)


def add_inline_brackets(s, x, y, w, h, size, color=INK, width_pt=1.0):
    """Two opposing brackets — left-open and right-open — to frame inline text.
       size is the leg length. ( ⌐  text  ⌎ )"""
    # Left bracket (⌐) — top horizontal + left vertical
    add_line(s, x, y, x + size, y, color, width_pt)
    add_line(s, x, y, x, y + h, color, width_pt)
    add_line(s, x, y + h, x + size, y + h, color, width_pt)
    # Right bracket (⌎) — top horizontal + right vertical
    add_line(s, x + w - size, y, x + w, y, color, width_pt)
    add_line(s, x + w, y, x + w, y + h, color, width_pt)
    add_line(s, x + w - size, y + h, x + w, y + h, color, width_pt)


def add_text(s, x, y, w, h, text, *,
             font=FONT_DISPLAY, size=14, bold=False, italic=False,
             color=INK, align="left", anchor="top", tracking=None, line_height=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    anchor_map = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]

    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        if line_height is not None:
            p.line_spacing = line_height
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if tracking is not None:
            # set character spacing (hundredths of a point in OOXML)
            r._r.get_or_add_rPr().set('spc', str(int(tracking * 100)))
    return tb


def add_circle(s, cx, cy, r, *, fill=None, line_color=INK, line_width=1.0, dash=None):
    """Add an oval centered at (cx, cy) with radius r."""
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    if fill is None:
        shp.fill.background()  # no fill
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line_color
    shp.line.width = Pt(line_width)
    if dash:
        ln = shp.line._get_or_add_ln()
        for el in ln.findall(qn('a:prstDash')):
            ln.remove(el)
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', dash)
    # Strip default shadow
    sp_pr = shp.fill._xPr
    return shp


def add_peak_icon(s, cx, cy, w, color=INK, width_pt=1.0):
    """Tiny mountain-peak chevron + horizontal line, matching reference icon."""
    half = w / 2
    # Chevron: two strokes meeting above center
    apex_x, apex_y = cx - w * 0.05, cy - w * 0.18
    add_line(s, cx - half * 0.85, cy + w * 0.05, apex_x, apex_y, color, width_pt)
    add_line(s, apex_x, apex_y, cx + half * 0.1, cy + w * 0.05, color, width_pt)
    # Horizontal "ground" tick on the right
    add_line(s, cx + half * 0.2, cy + w * 0.05, cx + half * 0.85, cy + w * 0.05, color, width_pt)


def add_bracket_button(s, x, y, w, h, label, *, size=Inches(0.16),
                       text_size=11, text_color=INK, bracket_color=INK,
                       width_pt=1.0, font=FONT_MONO, tracking=2.4):
    """A button styled as 4 corner brackets framing the label — no fill, no border.
       Returns the bracket leg size used (in EMU) so callers can chain hover effects."""
    add_corner_brackets(s, x, y, w, h, size, color=bracket_color, width_pt=width_pt)
    add_text(s, x, y, w, h, label,
             font=font, size=text_size, color=text_color,
             align="center", anchor="middle", tracking=tracking, bold=True)


def add_filled_bracket_button(s, x, y, w, h, label, *, size=Inches(0.16),
                              text_size=11, fill=INK, text_color=SAGE,
                              font=FONT_MONO, tracking=2.4):
    """Primary CTA variant — ink-filled rectangle with sage brackets on top."""
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    # corner brackets in inverse color, inset slightly
    inset = Emu(int(size * 0.15))
    add_corner_brackets(s, x + inset, y + inset, w - inset * 2, h - inset * 2,
                        size, color=text_color, width_pt=1.0)
    add_text(s, x, y, w, h, label,
             font=font, size=text_size, color=text_color,
             align="center", anchor="middle", tracking=tracking, bold=True)


# ─────────────────────────────────────────────────────────────────────────────
# REUSABLE SLIDE FURNITURE
# ─────────────────────────────────────────────────────────────────────────────

PAGE_MARGIN = Inches(0.55)
FRAME_BRACKET_SIZE = Inches(0.32)


def add_page_frame(s):
    """Outer corner brackets that frame every content slide."""
    add_corner_brackets(s, PAGE_MARGIN, PAGE_MARGIN,
                        SW - PAGE_MARGIN * 2, SH - PAGE_MARGIN * 2,
                        FRAME_BRACKET_SIZE, color=INK, width_pt=0.75)


def add_footer(s, page_num, total, section_label):
    """Tiny footer strip — page index left, section right."""
    add_text(s, Inches(0.8), SH - Inches(0.55), Inches(4), Inches(0.3),
             f"{page_num:02d} / {total:02d}",
             font=FONT_MONO, size=8.5, color=INK_MUTE,
             tracking=2.0, align="left", anchor="middle")
    add_text(s, SW - Inches(4.8), SH - Inches(0.55), Inches(4), Inches(0.3),
             section_label,
             font=FONT_MONO, size=8.5, color=INK_MUTE,
             tracking=2.4, align="right", anchor="middle")


def add_dotted_h(s, x, y, w, color=INK, width_pt=0.75):
    add_line(s, x, y, x + w, y, color=color, width_pt=width_pt, dash="sysDot")


def add_dotted_v(s, x, y, h, color=INK, width_pt=0.75):
    add_line(s, x, y, x, y + h, color=color, width_pt=width_pt, dash="sysDot")


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_page_frame(s)

# Centered circle with peak icon
cx, cy = SW / 2, Inches(3.05)
r_outer = Inches(0.85)
# Outer dashed ring (subtle, sits behind)
add_circle(s, cx + Inches(0.06), cy + Inches(0.04), r_outer + Inches(0.05),
           line_color=INK, line_width=0.6, dash="sysDot")
# Main solid ring
add_circle(s, cx, cy, r_outer, line_color=INK, line_width=1.0)
# Peak icon inside
add_peak_icon(s, cx, cy, Inches(0.85), color=INK, width_pt=1.0)

# Bracket-framed title
title = "PROJECT ATLAS"
title_w = Inches(4.4)
title_h = Inches(0.85)
title_x = (SW - title_w) / 2
title_y = Inches(4.45)
add_inline_brackets(s, title_x, title_y, title_w, title_h, Inches(0.22),
                    color=INK, width_pt=1.0)
add_text(s, title_x, title_y, title_w, title_h, title,
         font=FONT_DISPLAY, size=30, color=INK, bold=True,
         align="center", anchor="middle", tracking=8.0)

# Dotted underline tagline
tagline_y = title_y + title_h + Inches(0.45)
tagline_w = Inches(5.2)
add_dotted_h(s, (SW - tagline_w) / 2, tagline_y, tagline_w, width_pt=0.6)
add_text(s, (SW - tagline_w) / 2, tagline_y + Inches(0.12), tagline_w, Inches(0.3),
         "A QUIET CARTOGRAPHY OF WHAT COMES NEXT",
         font=FONT_MONO, size=9, color=INK_SOFT,
         align="center", anchor="top", tracking=3.4)

# Tiny corner metadata
add_text(s, PAGE_MARGIN + Inches(0.05), PAGE_MARGIN + Inches(0.05),
         Inches(3), Inches(0.3), "VOL. 01  ·  EDITION I",
         font=FONT_MONO, size=8, color=INK_MUTE, tracking=2.6)
add_text(s, SW - PAGE_MARGIN - Inches(3.05), PAGE_MARGIN + Inches(0.05),
         Inches(3), Inches(0.3), "2026  —  INTERNAL DRAFT",
         font=FONT_MONO, size=8, color=INK_MUTE, tracking=2.6, align="right")

add_footer(s, 1, 6, "COVER")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SECTION DIVIDER
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_page_frame(s)

# Big "01" numeral
num_x = Inches(1.4)
num_y = Inches(2.4)
add_text(s, num_x, num_y, Inches(3.5), Inches(2.6), "01",
         font=FONT_DISPLAY, size=200, color=INK, bold=False,
         align="left", anchor="top", tracking=-6.0)

# Circle marker after the numeral (echoes cover)
mark_cx = Inches(4.05)
mark_cy = Inches(3.55)
add_circle(s, mark_cx, mark_cy, Inches(0.16), line_color=INK, line_width=0.9)

# Dotted trail to the section title
trail_y = mark_cy
add_line(s, mark_cx + Inches(0.18), trail_y,
         Inches(6.8), trail_y, color=INK, width_pt=0.6, dash="sysDot")

# Section title in inline brackets
sec_label_x = Inches(6.95)
sec_label_y = Inches(3.22)
sec_label_w = Inches(5.4)
sec_label_h = Inches(0.7)
add_inline_brackets(s, sec_label_x, sec_label_y, sec_label_w, sec_label_h,
                    Inches(0.16), color=INK, width_pt=0.9)
add_text(s, sec_label_x, sec_label_y, sec_label_w, sec_label_h,
         "DISCOVERY",
         font=FONT_DISPLAY, size=26, color=INK, bold=True,
         align="center", anchor="middle", tracking=10.0)

# Subtitle below in mono
add_text(s, sec_label_x, sec_label_y + sec_label_h + Inches(0.12),
         sec_label_w, Inches(0.4),
         "PHASE ONE  ·  WEEK 01–04",
         font=FONT_MONO, size=10, color=INK_SOFT,
         align="center", anchor="top", tracking=3.6)

# Tiny tag in corner
add_text(s, PAGE_MARGIN + Inches(0.05), PAGE_MARGIN + Inches(0.05),
         Inches(4), Inches(0.3), "SECTION  ·  01 OF 04",
         font=FONT_MONO, size=8, color=INK_MUTE, tracking=2.6)

add_footer(s, 2, 6, "DISCOVERY")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTENT (bracketed nodes connected by dotted lines)
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_page_frame(s)

# Title
add_text(s, Inches(0.95), Inches(0.85), Inches(8), Inches(0.45),
         "THE METHOD",
         font=FONT_DISPLAY, size=22, color=INK, bold=True, tracking=4.0)
add_text(s, Inches(0.95), Inches(1.32), Inches(8), Inches(0.32),
         "Four passes, each feeding the next — no step skipped.",
         font=FONT_BODY, size=12, color=INK_SOFT, italic=True)

# Four bracketed cards arranged in a horizontal flow
card_w = Inches(2.55)
card_h = Inches(2.2)
card_y = Inches(2.55)
gap = Inches(0.35)
total_w = card_w * 4 + gap * 3
start_x = (SW - total_w) / 2

cards = [
    ("01", "OBSERVE",  "Walk the ground. Take only notes, not conclusions."),
    ("02", "ANNOTATE", "Mark every assumption — especially the comfortable ones."),
    ("03", "MODEL",    "Sketch the smallest version that still explains the whole."),
    ("04", "TEST",     "Break it cheaply before someone else breaks it loudly."),
]

for i, (num, label, body) in enumerate(cards):
    x = start_x + i * (card_w + gap)
    # corner brackets
    add_corner_brackets(s, x, card_y, card_w, card_h, Inches(0.2),
                        color=INK, width_pt=0.9)
    # numeral
    add_text(s, x + Inches(0.22), card_y + Inches(0.18), card_w - Inches(0.4), Inches(0.4),
             num, font=FONT_MONO, size=11, color=INK_MUTE, tracking=2.0, bold=True)
    # tiny circle marker top-right
    add_circle(s, x + card_w - Inches(0.28), card_y + Inches(0.3),
               Inches(0.07), line_color=INK, line_width=0.7)
    # label
    add_text(s, x + Inches(0.22), card_y + Inches(0.62), card_w - Inches(0.4), Inches(0.45),
             label, font=FONT_DISPLAY, size=18, color=INK, bold=True, tracking=3.2)
    # dotted divider
    add_dotted_h(s, x + Inches(0.22), card_y + Inches(1.16),
                 card_w - Inches(0.44), width_pt=0.55)
    # body copy
    add_text(s, x + Inches(0.22), card_y + Inches(1.28),
             card_w - Inches(0.44), card_h - Inches(1.4),
             body, font=FONT_BODY, size=11, color=INK_SOFT,
             align="left", anchor="top", line_height=1.35)

# Dotted connectors BETWEEN cards (across the gaps)
for i in range(3):
    x_from = start_x + (i + 1) * card_w + i * gap
    x_to   = x_from + gap
    y_mid  = card_y + card_h / 2
    add_dotted_h(s, x_from, y_mid, gap, width_pt=0.6)

add_footer(s, 3, 6, "DISCOVERY · METHOD")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — STATS (three big hairline numbers)
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_page_frame(s)

# Title
add_text(s, Inches(0.95), Inches(0.85), Inches(8), Inches(0.45),
         "WHAT WE FOUND",
         font=FONT_DISPLAY, size=22, color=INK, bold=True, tracking=4.0)
add_text(s, Inches(0.95), Inches(1.32), Inches(10), Inches(0.32),
         "Three numbers that reframed the brief.",
         font=FONT_BODY, size=12, color=INK_SOFT, italic=True)

stats = [
    ("87",  "%",  "of users abandoned\nbefore the third step"),
    ("4.2", "x",  "longer sessions when\nthe first prompt was open-ended"),
    ("01",  "",   "single feature drove\nhalf of all returns"),
]

stat_block_w = Inches(3.7)
gap = Inches(0.3)
total_w = stat_block_w * 3 + gap * 2
start_x = (SW - total_w) / 2
block_y = Inches(2.5)
block_h = Inches(3.7)

for i, (num, suffix, label) in enumerate(stats):
    x = start_x + i * (stat_block_w + gap)
    # Big hairline number — accent CLAY for the middle one only
    num_color = CLAY if i == 1 else INK
    add_text(s, x, block_y, stat_block_w, Inches(2.0),
             num, font=FONT_DISPLAY, size=130, color=num_color, bold=False,
             align="center", anchor="middle", tracking=-3.0)
    # Tiny suffix character to the right of the number
    if suffix:
        add_text(s, x, block_y + Inches(0.6), stat_block_w + Inches(1.6), Inches(0.6),
                 suffix, font=FONT_MONO, size=22, color=num_color, bold=True,
                 align="right", anchor="top")
    # Dotted underline
    underline_w = Inches(1.6)
    add_dotted_h(s, x + (stat_block_w - underline_w) / 2,
                 block_y + Inches(2.25), underline_w, width_pt=0.6)
    # Mini circle marker
    add_circle(s, x + stat_block_w / 2, block_y + Inches(2.45),
               Inches(0.07), line_color=INK, line_width=0.7)
    # Label
    add_text(s, x + Inches(0.1), block_y + Inches(2.65),
             stat_block_w - Inches(0.2), Inches(1.0),
             label, font=FONT_BODY, size=12, color=INK_SOFT,
             align="center", anchor="top", line_height=1.4)

add_footer(s, 4, 6, "FINDINGS")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TWO-COLUMN COMPARE
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_page_frame(s)

add_text(s, Inches(0.95), Inches(0.85), Inches(8), Inches(0.45),
         "BEFORE  ·  AFTER",
         font=FONT_DISPLAY, size=22, color=INK, bold=True, tracking=4.0)
add_text(s, Inches(0.95), Inches(1.32), Inches(10), Inches(0.32),
         "Same goal, fewer assumptions in the way.",
         font=FONT_BODY, size=12, color=INK_SOFT, italic=True)

col_w = Inches(5.5)
col_h = Inches(4.5)
col_y = Inches(2.25)
gap = Inches(0.6)
total_w = col_w * 2 + gap
start_x = (SW - total_w) / 2

# Left column — BEFORE (muted)
lx = start_x
add_corner_brackets(s, lx, col_y, col_w, col_h, Inches(0.24),
                    color=INK, width_pt=0.9)
add_text(s, lx + Inches(0.4), col_y + Inches(0.28), col_w - Inches(0.8), Inches(0.4),
         "BEFORE",
         font=FONT_MONO, size=10, color=INK_MUTE, tracking=3.6, bold=True)
add_text(s, lx + Inches(0.4), col_y + Inches(0.65), col_w - Inches(0.8), Inches(0.8),
         "Eleven required fields",
         font=FONT_DISPLAY, size=26, color=INK, bold=True, tracking=0)
add_dotted_h(s, lx + Inches(0.4), col_y + Inches(1.6), col_w - Inches(0.8), width_pt=0.55)

left_points = [
    "Three-page onboarding form",
    "Mandatory company size, role, sector",
    "Email confirmation before any value shown",
    "Hard-blocked behind a paywall on first launch",
]
for j, point in enumerate(left_points):
    yj = col_y + Inches(1.85) + Inches(0.55) * j
    # tiny bracket bullet (just a left ⌐)
    add_line(s, lx + Inches(0.4), yj + Inches(0.05),
             lx + Inches(0.52), yj + Inches(0.05), INK, 0.8)
    add_line(s, lx + Inches(0.4), yj + Inches(0.05),
             lx + Inches(0.4), yj + Inches(0.22), INK, 0.8)
    add_text(s, lx + Inches(0.65), yj, col_w - Inches(1.0), Inches(0.45),
             point, font=FONT_BODY, size=12.5, color=INK_SOFT, anchor="top")

# Right column — AFTER (CLAY accent label, otherwise same ink)
rx = start_x + col_w + gap
add_corner_brackets(s, rx, col_y, col_w, col_h, Inches(0.24),
                    color=INK, width_pt=0.9)
add_text(s, rx + Inches(0.4), col_y + Inches(0.28), col_w - Inches(0.8), Inches(0.4),
         "AFTER",
         font=FONT_MONO, size=10, color=CLAY, tracking=3.6, bold=True)
add_text(s, rx + Inches(0.4), col_y + Inches(0.65), col_w - Inches(0.8), Inches(0.8),
         "One question, one button",
         font=FONT_DISPLAY, size=26, color=INK, bold=True, tracking=0)
add_dotted_h(s, rx + Inches(0.4), col_y + Inches(1.6), col_w - Inches(0.8), width_pt=0.55)

right_points = [
    "Single open-ended prompt to start",
    "Everything else deferred to natural moments",
    "Email captured only when user asked to save",
    "First five sessions free, no card required",
]
for j, point in enumerate(right_points):
    yj = col_y + Inches(1.85) + Inches(0.55) * j
    add_line(s, rx + Inches(0.4), yj + Inches(0.05),
             rx + Inches(0.52), yj + Inches(0.05), INK, 0.8)
    add_line(s, rx + Inches(0.4), yj + Inches(0.05),
             rx + Inches(0.4), yj + Inches(0.22), INK, 0.8)
    add_text(s, rx + Inches(0.65), yj, col_w - Inches(1.0), Inches(0.45),
             point, font=FONT_BODY, size=12.5, color=INK, anchor="top")

# Dotted vertical divider between columns
mid_x = start_x + col_w + gap / 2
add_dotted_v(s, mid_x, col_y + Inches(0.35), col_h - Inches(0.7), width_pt=0.55)
# Tiny circle marker in middle of divider
add_circle(s, mid_x, col_y + col_h / 2, Inches(0.08),
           line_color=INK, line_width=0.7, fill=SAGE)

add_footer(s, 5, 6, "RESHAPE")

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CLOSING / CTA (with bracket-styled buttons)
# ═════════════════════════════════════════════════════════════════════════════
s = add_slide()
add_page_frame(s)

# Empty circle (mirrors cover, "completed")
cx, cy = SW / 2, Inches(2.5)
add_circle(s, cx, cy, Inches(0.6), line_color=INK, line_width=1.0, dash="sysDot")
# A tiny solid dot at center (the "you-are-here")
add_circle(s, cx, cy, Inches(0.07), fill=INK, line_color=INK, line_width=0.5)

# THANK YOU in brackets
ty_w = Inches(5.2)
ty_h = Inches(0.95)
ty_x = (SW - ty_w) / 2
ty_y = Inches(3.55)
add_inline_brackets(s, ty_x, ty_y, ty_w, ty_h, Inches(0.24),
                    color=INK, width_pt=1.0)
add_text(s, ty_x, ty_y, ty_w, ty_h, "THANK YOU",
         font=FONT_DISPLAY, size=34, color=INK, bold=True,
         align="center", anchor="middle", tracking=10.0)

# Dotted line and subline
sub_w = Inches(6.5)
add_dotted_h(s, (SW - sub_w) / 2, ty_y + ty_h + Inches(0.32), sub_w, width_pt=0.6)
add_text(s, (SW - sub_w) / 2, ty_y + ty_h + Inches(0.44), sub_w, Inches(0.3),
         "QUESTIONS, OBJECTIONS, BETTER IDEAS — ALL WELCOME",
         font=FONT_MONO, size=9, color=INK_SOFT,
         align="center", tracking=3.0)

# ── BUTTON ROW — bracket-styled ────────────────────────────────────────────
btn_y = Inches(5.7)
btn_h = Inches(0.55)
btn1_w = Inches(2.2)
btn2_w = Inches(2.4)
btn3_w = Inches(2.2)
btn_gap = Inches(0.3)
total_btn_w = btn1_w + btn2_w + btn3_w + btn_gap * 2
btn_x_start = (SW - total_btn_w) / 2

# Button 1 — secondary (outline brackets only)
add_bracket_button(s, btn_x_start, btn_y, btn1_w, btn_h, "READ THE DOCS",
                   size=Inches(0.13), text_size=10, width_pt=0.9)

# Button 2 — primary (filled ink, sage brackets)
add_filled_bracket_button(s, btn_x_start + btn1_w + btn_gap, btn_y,
                          btn2_w, btn_h, "GET STARTED",
                          size=Inches(0.13), text_size=10.5)

# Button 3 — secondary (outline)
add_bracket_button(s, btn_x_start + btn1_w + btn_gap + btn2_w + btn_gap,
                   btn_y, btn3_w, btn_h, "SAY HELLO",
                   size=Inches(0.13), text_size=10, width_pt=0.9)

# Tiny under-button hint
add_text(s, 0, btn_y + btn_h + Inches(0.18), SW, Inches(0.25),
         "atlas.studio  ·  hello@atlas.studio",
         font=FONT_MONO, size=8.5, color=INK_MUTE,
         align="center", tracking=2.6)

add_footer(s, 6, 6, "CLOSING")

# ─────────────────────────────────────────────────────────────────────────────
out = "outputs/sage-brackets-template/sage-brackets-template.pptx"
prs.save(out)
print(f"Wrote {out}")
